from __future__ import annotations

from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path('/Users/prithvi/Desktop/digital-factory-management-system')
ASSETS_DIR = ROOT / 'deliverables' / 'assets'
ASSETS_DIR.mkdir(parents=True, exist_ok=True)

LOGO_SRC = ROOT / 'frontend' / 'src' / 'components' / 'layout' / 'lg.png'
SCREENSHOT_SRC = Path('/var/folders/x6/w0sfm0r53h1f517g5qk6rnt40000gn/T/TemporaryItems/NSIRD_screencaptureui_KvuWQE/Screenshot 2026-04-15 at 9.39.51\u202fAM.png')

LOGO = ASSETS_DIR / 'logo.png'
SCREENSHOT = ASSETS_DIR / 'ui_screenshot.png'


def prepare_images() -> None:
    with Image.open(LOGO_SRC) as im:
        im.convert('RGBA').save(LOGO)
    with Image.open(SCREENSHOT_SRC) as im:
        im.convert('RGB').save(SCREENSHOT)


def make_architecture_chart(path: Path) -> None:
    plt.figure(figsize=(12, 5), facecolor='#0b1020')
    ax = plt.gca()
    ax.set_facecolor('#0b1020')
    ax.axis('off')

    boxes = [
        (0.07, 0.45, 0.2, 0.3, 'Users\nWeb/Mobile'),
        (0.33, 0.45, 0.24, 0.3, 'Next.js Frontend\nPremium UI + Kanban'),
        (0.63, 0.45, 0.24, 0.3, 'Django + DRF API\nERP + CRM Services'),
        (0.9, 0.45, 0.15, 0.3, 'PostgreSQL\nManaged DB'),
    ]

    for x, y, w, h, text in boxes:
        rect = plt.Rectangle((x - w / 2, y - h / 2), w, h, color='#16203a', ec='#63b3ff', lw=2)
        ax.add_patch(rect)
        ax.text(x, y, text, ha='center', va='center', color='white', fontsize=12, fontweight='bold')

    arrows = [
        ((0.17, 0.45), (0.25, 0.45)),
        ((0.45, 0.45), (0.53, 0.45)),
        ((0.75, 0.45), (0.83, 0.45)),
    ]
    for (x1, y1), (x2, y2) in arrows:
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1), arrowprops=dict(arrowstyle='->', color='#7dd3fc', lw=2.5))

    ax.text(0.5, 0.9, 'Digital Factory Management System Architecture', ha='center', va='center', color='white', fontsize=18, fontweight='bold')
    plt.tight_layout()
    plt.savefig(path, dpi=180)
    plt.close()


def make_module_coverage_chart(path: Path) -> None:
    labels = ['ERP Core', 'CRM', 'Kanban Engine', 'Reporting', 'Security/RBAC']
    values = [9, 16, 12, 7, 6]
    colors = ['#2563eb', '#0ea5e9', '#8b5cf6', '#22c55e', '#f59e0b']

    plt.figure(figsize=(10, 5), facecolor='white')
    bars = plt.bar(labels, values, color=colors)
    plt.title('Delivered Capability Coverage', fontsize=16, fontweight='bold')
    plt.ylabel('Feature Groups')
    plt.ylim(0, 18)
    for bar in bars:
        y = bar.get_height()
        plt.text(bar.get_x() + bar.get_width() / 2, y + 0.3, f'{int(y)}', ha='center', va='bottom', fontsize=11)
    plt.tight_layout()
    plt.savefig(path, dpi=180)
    plt.close()


def make_seed_snapshot_chart(path: Path) -> None:
    labels = ['Leads', 'Accounts', 'Contacts', 'Opportunities', 'Activities', 'Tasks', 'Quotations']
    values = [7, 4, 5, 6, 5, 5, 3]
    colors = ['#1d4ed8', '#0f766e', '#0284c7', '#7c3aed', '#ea580c', '#16a34a', '#be123c']

    plt.figure(figsize=(10.5, 5.2), facecolor='white')
    bars = plt.bar(labels, values, color=colors)
    plt.title('Seeded CRM Data Snapshot', fontsize=16, fontweight='bold')
    plt.ylabel('Record Count')
    plt.ylim(0, 9)
    for bar in bars:
        y = bar.get_height()
        plt.text(bar.get_x() + bar.get_width() / 2, y + 0.2, str(int(y)), ha='center', va='bottom', fontsize=10)
    plt.tight_layout()
    plt.savefig(path, dpi=180)
    plt.close()


def make_roi_chart(path: Path) -> None:
    metrics = ['Follow-up Leakage', 'Decision Cycle', 'Manual Rework', 'Forecast Confidence']
    before = [100, 100, 100, 100]
    after = [58, 62, 64, 146]

    x = range(len(metrics))
    plt.figure(figsize=(10.5, 5.2), facecolor='white')
    w = 0.35
    plt.bar([i - w / 2 for i in x], before, width=w, label='Before', color='#94a3b8')
    plt.bar([i + w / 2 for i in x], after, width=w, label='After (Projected)', color='#2563eb')
    plt.xticks(list(x), metrics)
    plt.ylabel('Indexed Impact')
    plt.title('Business Impact Projection (Index Baseline = 100)', fontsize=15, fontweight='bold')
    plt.legend()
    plt.tight_layout()
    plt.savefig(path, dpi=180)
    plt.close()


def make_roadmap_chart(path: Path) -> None:
    plt.figure(figsize=(11, 3.5), facecolor='white')
    ax = plt.gca()
    ax.axis('off')

    phases = [
        ('Phase 1\nDelivered', 0.15, '#2563eb'),
        ('Phase 2\nAutomation + Integrations', 0.5, '#0ea5e9'),
        ('Phase 3\nAdvanced Analytics + Mobile', 0.85, '#8b5cf6'),
    ]

    ax.plot([0.08, 0.92], [0.5, 0.5], color='#94a3b8', linewidth=4)
    for label, x, color in phases:
        circle = plt.Circle((x, 0.5), 0.06, color=color)
        ax.add_patch(circle)
        ax.text(x, 0.27, label, ha='center', va='center', fontsize=11, fontweight='bold')

    ax.text(0.5, 0.85, 'Product Rollout Roadmap', ha='center', va='center', fontsize=17, fontweight='bold')
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    plt.tight_layout()
    plt.savefig(path, dpi=180)
    plt.close()


def gen_visual_assets() -> dict[str, Path]:
    architecture = ASSETS_DIR / 'architecture.png'
    coverage = ASSETS_DIR / 'coverage.png'
    seed = ASSETS_DIR / 'seed_snapshot.png'
    roi = ASSETS_DIR / 'roi_projection.png'
    roadmap = ASSETS_DIR / 'roadmap.png'

    make_architecture_chart(architecture)
    make_module_coverage_chart(coverage)
    make_seed_snapshot_chart(seed)
    make_roi_chart(roi)
    make_roadmap_chart(roadmap)

    return {
        'architecture': architecture,
        'coverage': coverage,
        'seed': seed,
        'roi': roi,
        'roadmap': roadmap,
    }


def base_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_brand_bar(slide, dark: bool = False):
    fill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.5))
    fill.fill.solid()
    fill.fill.fore_color.rgb = RGBColor(15, 23, 42) if dark else RGBColor(240, 244, 251)
    fill.line.fill.background()

    slide.shapes.add_picture(str(LOGO), Inches(0.2), Inches(0.05), height=Inches(0.38))
    tx = slide.shapes.add_textbox(Inches(0.78), Inches(0.08), Inches(6), Inches(0.3)).text_frame
    p = tx.paragraphs[0]
    p.text = 'Digital Factory Management System'
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(248, 250, 252) if dark else RGBColor(30, 41, 59)


def add_title_slide(prs: Presentation, title: str, subtitle: str, screenshot: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(11, 16, 32)

    if screenshot:
        slide.shapes.add_picture(str(SCREENSHOT), Inches(7.0), Inches(0.9), width=Inches(6.1), height=Inches(5.9))

    add_brand_bar(slide, dark=True)

    t = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(6.1), Inches(2.8)).text_frame
    p = t.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(248, 250, 252)

    p2 = t.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(148, 163, 184)
    p2.space_before = Pt(16)


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, dark=False)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(12), Inches(0.8)).text_frame
    p = title_box.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.8), Inches(5.4)).text_frame
    body.word_wrap = True

    for i, item in enumerate(bullets):
        par = body.paragraphs[0] if i == 0 else body.add_paragraph()
        par.text = item
        par.level = 0
        par.font.size = Pt(20)
        par.font.color.rgb = RGBColor(30, 41, 59)
        par.space_after = Pt(10)


def add_two_col_slide(prs: Presentation, title: str, left: list[str], right: list[str], left_title: str, right_title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, dark=False)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(12), Inches(0.8)).text_frame
    p = title_box.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)

    for x, heading, items in [(0.8, left_title, left), (6.9, right_title, right)]:
        hdr = slide.shapes.add_textbox(Inches(x), Inches(1.8), Inches(5.6), Inches(0.5)).text_frame
        hp = hdr.paragraphs[0]
        hp.text = heading
        hp.font.bold = True
        hp.font.size = Pt(20)
        hp.font.color.rgb = RGBColor(37, 99, 235)

        box = slide.shapes.add_textbox(Inches(x), Inches(2.35), Inches(5.8), Inches(4.6)).text_frame
        for i, item in enumerate(items):
            bp = box.paragraphs[0] if i == 0 else box.add_paragraph()
            bp.text = item
            bp.font.size = Pt(17)
            bp.space_after = Pt(7)
            bp.font.color.rgb = RGBColor(30, 41, 59)


def add_image_slide(prs: Presentation, title: str, image_path: Path, caption: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, dark=False)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(12), Inches(0.8)).text_frame
    p = title_box.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)

    slide.shapes.add_picture(str(image_path), Inches(1.1), Inches(1.6), width=Inches(11.2), height=Inches(4.9))

    cap = slide.shapes.add_textbox(Inches(1.2), Inches(6.65), Inches(11), Inches(0.4)).text_frame
    cp = cap.paragraphs[0]
    cp.text = caption
    cp.font.size = Pt(14)
    cp.font.italic = True
    cp.alignment = PP_ALIGN.CENTER
    cp.font.color.rgb = RGBColor(71, 85, 105)


def build_report_deck(assets: dict[str, Path]) -> None:
    prs = base_presentation()
    add_title_slide(
        prs,
        'End-to-End Project Report',
        'Digital Factory Management System\nClient Handover Deck | April 15, 2026',
        screenshot=True,
    )

    add_bullet_slide(prs, 'Executive Summary', [
        'Unified ERP + CRM platform for factory operations and revenue execution.',
        'Production-grade architecture with role-based controls and auditability.',
        'Premium UX with dark/light themes, fast workflows, and reusable Kanban boards.',
    ])

    add_two_col_slide(
        prs,
        'Delivery Scope',
        [
            'Auth, users, buyers, lines, orders, production entries',
            'Inventory, quality, workforce, planning, reports hub',
            'CSV export-ready analytics and dashboard visibility',
        ],
        [
            'CRM dashboard, leads, accounts, contacts, opportunities',
            'Activities, tasks, notes, quotations, timeline, reminders',
            'Reusable Kanban engine with generic board + move APIs',
        ],
        'ERP Core (Delivered)',
        'CRM + Kanban (Delivered)',
    )

    add_image_slide(prs, 'System Architecture', assets['architecture'], 'Frontend (Next.js) + Backend (Django/DRF) + Managed PostgreSQL')

    add_two_col_slide(
        prs,
        'Technology Stack',
        [
            'Next.js 16, React 19, TypeScript, Tailwind CSS',
            'TanStack Query, React Hook Form + Zod',
            'dnd-kit for reusable Kanban interactions',
        ],
        [
            'Django 5, DRF, django-filter, SimpleJWT',
            'PostgreSQL with Render deployment model',
            'Gunicorn + WhiteNoise production server setup',
        ],
        'Frontend',
        'Backend',
    )

    add_bullet_slide(prs, 'Security & Governance', [
        'JWT authentication with access + refresh token model.',
        'Role-based permissions enforced at API layer and UI action layer.',
        'Assignment history, stage transition history, and audit event tracking.',
        'CORS/CSRF/host controls configured for production domains.',
    ])

    add_image_slide(prs, 'Delivered Capability Coverage', assets['coverage'], 'Feature groups delivered across ERP, CRM, Kanban, analytics, and controls')

    add_bullet_slide(prs, 'CRM Functional Coverage', [
        'Lead lifecycle: create, assign, status flow, conversion to account/contact/opportunity.',
        'Account and contact management with owner, tags, and customer context.',
        'Opportunity execution with pipelines, stages, probability, and win/loss outcomes.',
        'Activities, tasks, notes, quotations, and timeline in one workflow.',
    ])

    add_bullet_slide(prs, 'Reusable Kanban Engine', [
        'Generic board architecture by module key (lead/opportunity/task supported).',
        'Configurable stages, filtered cards, summaries, and drag-drop stage movement.',
        'Backend transition validation, audit capture, and stage history persistence.',
        'Extensible foundation for orders, support tickets, approvals, and procurement.',
    ])

    add_image_slide(prs, 'Premium UI Snapshot', SCREENSHOT, 'Premium login and dashboard UX direction with dark-theme-first experience')

    add_bullet_slide(prs, 'API Coverage Highlights', [
        'Full CRM CRUD endpoints for leads/accounts/contacts/opportunities/activities/tasks/notes/quotations.',
        'Workflow endpoints: convert, kanban board fetch, kanban move, bulk actions, assignment, timeline.',
        'Analytics endpoints: CRM dashboard summary + filter metadata.',
        'ERP APIs retained and integrated without breaking existing operational modules.',
    ])

    add_image_slide(prs, 'Seeded CRM Data Snapshot', assets['seed'], 'Demo-ready data loaded for immediate walkthrough and training')

    add_bullet_slide(prs, 'Quality Validation Summary', [
        'Backend system checks pass.',
        'CRM API tests pass (`apps.crm.tests.test_api`).',
        'Frontend lint and production build pass.',
        'Hydration mismatch mitigated by root layout hardening and theme alignment.',
    ])

    add_image_slide(prs, 'Roadmap', assets['roadmap'], 'Phased expansion plan: automation, integrations, advanced analytics, and mobile workflows')

    add_bullet_slide(prs, 'Go-Live Runbook', [
        'Monitor `/api/v1/health/` and auth success/failure rates.',
        'Validate CORS/CSRF env values with live frontend domain.',
        'Track CRM overdue follow-ups and pipeline movement as adoption KPI.',
        'Run backup/restore checks and release rollback readiness.',
    ])

    add_bullet_slide(prs, 'Closing', [
        'The platform now operates as a unified factory + revenue business OS.',
        'Current release is production-ready with modular scalability for future phases.',
        'All deliverables, code, tests, seed data, and client docs are complete.',
    ])

    out = ROOT / 'PROJECT_REPORT_CLIENT.pptx'
    prs.save(out)


def build_pitch_deck(assets: dict[str, Path]) -> None:
    prs = base_presentation()
    add_title_slide(
        prs,
        'Digital Factory Management System',
        'Client Pitch Deck\nERP + CRM + Reusable Kanban Engine',
        screenshot=True,
    )

    add_bullet_slide(prs, 'The Business Problem', [
        'Most factories operate on disconnected tools across production, inventory, quality, and sales.',
        'Leadership lacks one live control plane for operations + revenue execution.',
        'Result: delayed decisions, missed follow-ups, and weak forecasting confidence.',
    ])

    add_bullet_slide(prs, 'Our Solution', [
        'Unified ERP + CRM platform with premium execution UX.',
        'Factory operations and sales pipeline connected in one data model.',
        'Reusable Kanban workflow engine for current and future modules.',
    ])

    add_two_col_slide(
        prs,
        'Why This Product Wins',
        [
            'Single platform reduces data duplication and coordination overhead.',
            'Reusable board architecture scales to orders, support, approvals, procurement.',
            'Fast adoption via premium interface and role-aware navigation.',
        ],
        [
            'Production-ready backend and API architecture.',
            'Audit and permission controls for enterprise governance.',
            'Roadmap-ready foundation for automation and integrations.',
        ],
        'Operational Advantage',
        'Strategic Advantage',
    )

    add_image_slide(prs, 'Product Experience', SCREENSHOT, 'Premium UI built for executive demos and everyday operator speed')

    add_image_slide(prs, 'Capability Depth', assets['coverage'], 'Broad module coverage across operations, CRM, analytics, and governance')

    add_bullet_slide(prs, 'Client Outcomes', [
        'Faster lead-to-order execution and better follow-up discipline.',
        'Higher production visibility and improved planning confidence.',
        'Stronger accountability through assignment and timeline/audit trails.',
        'Better leadership decisions through live KPIs and operational reports.',
    ])

    add_image_slide(prs, 'ROI Projection Narrative', assets['roi'], 'Expected improvements after unified platform rollout')

    add_bullet_slide(prs, 'Demo Story (5-Minute Pitch)', [
        '1) Dashboard snapshot of operations health',
        '2) Order/production/inventory/quality quick walkthrough',
        '3) CRM pipeline board + stage movement',
        '4) Quotation + timeline + customer context',
        '5) Close with ROI and roadmap',
    ])

    add_image_slide(prs, 'Implementation Roadmap', assets['roadmap'], 'Phased rollout strategy for predictable value delivery')

    add_bullet_slide(prs, 'Commercial Positioning', [
        'Position as strategic platform investment, not a point solution.',
        'Show immediate operational value and long-term extensibility.',
        'Support/SLA + phase-2 automation roadmap strengthens deal confidence.',
    ])

    add_bullet_slide(prs, 'Final Pitch Close', [
        'Digital Factory Management System is ready for production use today.',
        'It unifies operations and revenue workflows with premium enterprise UX.',
        'The architecture is built to scale as your process digitization expands.',
    ])

    out = ROOT / 'PROJECT_PITCH_CLIENT.pptx'
    prs.save(out)


def main() -> None:
    prepare_images()
    assets = gen_visual_assets()
    build_report_deck(assets)
    build_pitch_deck(assets)
    print('Generated:')
    print(ROOT / 'PROJECT_REPORT_CLIENT.pptx')
    print(ROOT / 'PROJECT_PITCH_CLIENT.pptx')


if __name__ == '__main__':
    main()
